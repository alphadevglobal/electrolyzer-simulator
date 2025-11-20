"""
Cria apresentação PPTX profissional e lúdica
Para o Projeto IA AV3
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os


def create_title_slide(prs, title, subtitle):
    """Cria slide de título"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 153)

    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)

    return slide


def create_content_slide(prs, title, content_points):
    """Cria slide de conteúdo com bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 153)

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    for point in content_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(20)
        p.space_before = Pt(10)

    return slide


def create_table_slide(prs, title, headers, rows):
    """Cria slide com tabela"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Adiciona tabela
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(4)

    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table

    # Header
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 153)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Rows
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)

    return slide


def create_conclusion_slide(prs, title, points, highlight):
    """Cria slide de conclusão com destaque"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 51)

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.space_before = Pt(10)

    # Adiciona caixa de destaque
    left = Inches(1.5)
    top = Inches(5)
    width = Inches(7)
    height = Inches(1.5)

    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.text = highlight
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 102, 51)

    return slide


def main():
    """Cria apresentação completa"""
    print("🎨 Criando apresentação PPTX...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Capa
    create_title_slide(
        prs,
        "Projeto IA AV3\nClassificação de Energia",
        "Inteligência Artificial Computacional\nProf. Ms. Cynthia Moreira Maia\nMateus Gomes Macário"
    )

    # Slide 2: Agenda
    create_content_slide(
        prs,
        "📋 Agenda",
        [
            "1. Introdução e Objetivo",
            "2. Dataset Escolhido",
            "3. Algoritmos Implementados",
            "4. Metodologia e Validação",
            "5. Resultados Obtidos",
            "6. Análise Comparativa",
            "7. Conclusões",
            "8. Demonstração Prática"
        ]
    )

    # Slide 3: Introdução
    create_content_slide(
        prs,
        "🎯 Introdução",
        [
            "Objetivo: Avaliar algoritmos de classificação",
            "Desafio: Implementação 100% MANUAL",
            "❌ SEM pandas",
            "❌ SEM scikit-learn",
            "✅ Tudo do ZERO!",
            "Problema: Classificação de consumo de energia"
        ]
    )

    # Slide 4: Dataset
    create_content_slide(
        prs,
        "📊 Dataset: Appliances Energy",
        [
            "Fonte: OpenML (ID: 46283)",
            "📈 Instâncias: 5,000 amostras",
            "🔢 Features: 28 atributos",
            "🎯 Target: Consumo de energia (binarizado)",
            "✅ Atende TODOS os requisitos (>10 atributos, >1000 instâncias)",
            "🔗 Alinhado com pesquisa em eficiência energética"
        ]
    )

    # Slide 5: Distribuição
    create_content_slide(
        prs,
        "⚖️ Distribuição dos Dados",
        [
            "Classes Balanceadas:",
            "   • Classe 0 (Baixo consumo): 2,500 amostras (50%)",
            "   • Classe 1 (Alto consumo): 2,500 amostras (50%)",
            "",
            "Pré-processamento aplicado:",
            "   • Normalização Z-Score (implementada manualmente)",
            "   • Binarização pela mediana",
            "   • Divisão estratificada K-Folds"
        ]
    )

    # Slide 6: KNN
    create_content_slide(
        prs,
        "🔍 Algoritmo 1: K-Nearest Neighbors",
        [
            "Implementação Manual:",
            "   • Distância Euclidiana: √Σ(x₁-x₂)²",
            "   • Distância Manhattan: Σ|x₁-x₂|",
            "   • k = 5 vizinhos",
            "",
            "Características:",
            "   • Lazy learning (sem treino)",
            "   • Busca exaustiva",
            "   • Votação majoritária"
        ]
    )

    # Slide 7: Perceptron
    create_content_slide(
        prs,
        "🧠 Algoritmo 2: Perceptron",
        [
            "Implementação Manual:",
            "   • Taxa de aprendizado: 0.01",
            "   • Épocas: 50",
            "   • Estratégia: One-vs-Rest para multiclasse",
            "",
            "Regra de Atualização:",
            "   • w = w + η × erro × x",
            "   • b = b + η × erro",
            "",
            "Função de Ativação: Degrau (step)"
        ]
    )

    # Slide 8: MLP
    create_content_slide(
        prs,
        "🚀 Algoritmo 3: Multi-Layer Perceptron",
        [
            "Arquitetura da Rede:",
            "   • Input: 28 neurônios",
            "   • Hidden 1: 32 neurônios (ReLU)",
            "   • Hidden 2: 16 neurônios (ReLU)",
            "   • Output: 2 neurônios (Softmax)",
            "",
            "Treinamento:",
            "   • BACKPROPAGATION implementado manualmente!",
            "   • Mini-batch Gradient Descent (batch=64)",
            "   • Learning rate: 0.01, Epochs: 50"
        ]
    )

    # Slide 9: Naive Bayes
    create_content_slide(
        prs,
        "📊 Algoritmo 4: Naive Bayes",
        [
            "Duas Implementações:",
            "",
            "1️⃣ Versão Univariada:",
            "   • Assume independência entre features",
            "   • Mais rápida e simples",
            "",
            "2️⃣ Versão Multivariada:",
            "   • Considera correlação entre features",
            "   • Distribuição Gaussiana multivariada",
            "   • PDF implementada manualmente"
        ]
    )

    # Slide 10: Validação Cruzada
    create_content_slide(
        prs,
        "🔄 Validação Cruzada",
        [
            "K-Fold Cross-Validation (k=5)",
            "✅ Implementação 100% MANUAL!",
            "",
            "Processo:",
            "   1. Divide dados em 5 partes (folds)",
            "   2. Para cada fold:",
            "      • Treina com 4 folds",
            "      • Testa com 1 fold",
            "   3. Calcula média e desvio padrão",
            "",
            "Estratificação: Mantém proporção das classes"
        ]
    )

    # Slide 11: Métricas
    create_content_slide(
        prs,
        "📏 Métricas de Avaliação",
        [
            "Todas implementadas MANUALMENTE:",
            "",
            "1️⃣ Acurácia (Accuracy)",
            "   • Porcentagem de acertos",
            "",
            "2️⃣ Precisão (Precision)",
            "   • TP / (TP + FP)",
            "",
            "3️⃣ F1-Score",
            "   • 2 × (Precision × Recall) / (Precision + Recall)",
            "",
            "⏱️ Tempo de Treino e Teste"
        ]
    )

    # Slide 12: Resultados - Tabela
    create_table_slide(
        prs,
        "🏆 Resultados Completos",
        ["Classificador", "Acurácia", "F1-Score", "Tempo (s)"],
        [
            ["KNN (Euclidiana)", "80.04%", "80.05%", "13.09"],
            ["KNN (Manhattan)", "80.88%", "80.89%", "10.19"],
            ["Perceptron", "77.08%", "80.43%", "0.37"],
            ["MLP 🏆", "92.52%", "92.53%", "0.22"],
            ["NB (Univariado)", "90.90%", "90.90%", "0.12"],
            ["NB (Multivariado)", "91.62%", "91.62%", "0.06"]
        ]
    )

    # Slide 13: Análise por Desempenho
    create_content_slide(
        prs,
        "📊 Análise por Desempenho",
        [
            "Ranking de Acurácia:",
            "",
            "🥇 MLP: 92.52% ⭐",
            "🥈 Naive Bayes (Multivariado): 91.62%",
            "🥉 Naive Bayes (Univariado): 90.90%",
            "4️⃣ KNN (Manhattan): 80.88%",
            "5️⃣ KNN (Euclidiana): 80.04%",
            "6️⃣ Perceptron: 77.08%"
        ]
    )

    # Slide 14: Trade-off
    create_content_slide(
        prs,
        "⚖️ Trade-off: Desempenho vs Eficiência",
        [
            "Score de Eficiência = Acurácia / Tempo Total",
            "",
            "🥇 Naive Bayes (Multivariado): 14.08",
            "🥈 Naive Bayes (Univariado): 7.07",
            "🥉 MLP: 3.97",
            "4️⃣ Perceptron: 2.05",
            "5️⃣ KNN (Manhattan): 0.08",
            "6️⃣ KNN (Euclidiana): 0.06",
            "",
            "💡 NB Multivariado: Melhor custo-benefício!"
        ]
    )

    # Slide 15: Pontos Fortes e Limitações
    create_content_slide(
        prs,
        "💪 Pontos Fortes e Limitações",
        [
            "MLP:",
            "   ✅ Melhor acurácia (92.52%)",
            "   ✅ Captura relações não-lineares",
            "   ❌ Difícil de interpretar",
            "",
            "Naive Bayes:",
            "   ✅ Extremamente rápido",
            "   ✅ Alta acurácia (>90%)",
            "   ❌ Assume independência",
            "",
            "KNN:",
            "   ✅ Simples de implementar",
            "   ❌ Muito lento no teste"
        ]
    )

    # Slide 16: Insights
    create_content_slide(
        prs,
        "💡 Insights Obtidos",
        [
            "1️⃣ Deep Learning funciona!",
            "   • MLP superou todos os modelos clássicos",
            "",
            "2️⃣ Simplicidade pode ser melhor",
            "   • Naive Bayes: simples, rápido e eficaz",
            "",
            "3️⃣ KNN tem limitações de escala",
            "   • Impraticável para produção em larga escala",
            "",
            "4️⃣ Perceptron é limitado",
            "   • Problema não é linearmente separável",
            "   • MLP essencial para melhor desempenho"
        ]
    )

    # Slide 17: Requisitos Atendidos
    create_content_slide(
        prs,
        "✅ Requisitos Atendidos",
        [
            "✅ Dataset com >10 atributos (28 features)",
            "✅ Dataset com >1000 instâncias (5,000 samples)",
            "✅ KNN (Euclidiana + Manhattan)",
            "✅ Perceptron implementado",
            "✅ MLP com backpropagation manual",
            "✅ Naive Bayes (Univariado + Multivariado)",
            "✅ SEM pandas",
            "✅ SEM scikit-learn",
            "✅ Validação cruzada manual",
            "✅ Código fonte completo"
        ]
    )

    # Slide 18: Extras Implementados
    create_content_slide(
        prs,
        "🌟 Extras Implementados",
        [
            "Além dos requisitos:",
            "",
            "✨ GitHub Actions (CI/CD automático)",
            "☁️ Google Colab (notebook interativo)",
            "🚀 Deploy AWS Lambda (Free Tier)",
            "📊 Visualizações automáticas",
            "📚 Documentação completa",
            "🔗 Integração Web ↔ Colab",
            "📈 Gráficos comparativos",
            "💾 Exportação em múltiplos formatos"
        ]
    )

    # Slide 19: Tecnologias
    create_content_slide(
        prs,
        "🛠️ Tecnologias Utilizadas",
        [
            "Permitidas:",
            "   • Python 3.10",
            "   • NumPy (apenas para arrays)",
            "   • Matplotlib (apenas visualização)",
            "",
            "Implementado Manualmente:",
            "   • Todos os algoritmos de ML",
            "   • Carregamento de CSV",
            "   • Normalização",
            "   • Validação cruzada",
            "   • Todas as métricas"
        ]
    )

    # Slide 20: Demonstração
    create_content_slide(
        prs,
        "🎬 Demonstração Prática",
        [
            "Onde testar:",
            "",
            "1️⃣ Localmente:",
            "   cd projeto_ia_av3",
            "   python src/main.py",
            "",
            "2️⃣ Google Colab:",
            "   Notebook interativo com todos experimentos",
            "",
            "3️⃣ AWS Lambda:",
            "   API serverless para predições em tempo real",
            "",
            "4️⃣ GitHub Actions:",
            "   CI/CD automático a cada push"
        ]
    )

    # Slide 21: Conclusões
    create_conclusion_slide(
        prs,
        "🎯 Conclusões",
        [
            "✅ Todos os algoritmos implementados com sucesso",
            "✅ MLP: melhor performance (92.52%)",
            "✅ Naive Bayes: melhor eficiência",
            "✅ Validação robusta com K-Fold",
            "✅ Código 100% original e documentado",
            "✅ Infraestrutura profissional (CI/CD)"
        ],
        "🏆 Projeto Concluído com Excelência!"
    )

    # Slide 22: Aplicações Futuras
    create_content_slide(
        prs,
        "🚀 Aplicações Futuras",
        [
            "Para Produção:",
            "   • Naive Bayes para sistemas em tempo real",
            "   • MLP quando precisão é crítica",
            "",
            "Melhorias Possíveis:",
            "   • Testar com dataset completo (19,735 amostras)",
            "   • Ajustar hiperparâmetros",
            "   • Implementar ensemble methods",
            "   • Otimização de features",
            "",
            "Integração:",
            "   • Aplicação web React",
            "   • Mobile app",
            "   • Dashboard de monitoramento"
        ]
    )

    # Slide 23: Agradecimentos
    create_content_slide(
        prs,
        "🙏 Agradecimentos",
        [
            "Prof. Ms. Cynthia Moreira Maia",
            "   • Orientação na disciplina",
            "",
            "Karen Moura Fernandes",
            "   • Pesquisadora colaboradora",
            "",
            "Prof. Paulo Henrique Pereira Silva",
            "   • Orientação acadêmica",
            "",
            "OpenML",
            "   • Disponibilização do dataset",
            "",
            "Comunidade Open Source",
            "   • Inspiração e referências"
        ]
    )

    # Slide 24: Perguntas
    create_title_slide(
        prs,
        "❓ Perguntas?",
        "Estou à disposição para esclarecer dúvidas!\n\nMateus Gomes Macário\nEngenharia da Computação - UNIFOR"
    )

    # Salva apresentação
    output_path = "slides/APRESENTACAO_AV3.pptx"
    os.makedirs("slides", exist_ok=True)
    prs.save(output_path)

    print(f"✅ Apresentação criada: {output_path}")
    print(f"📊 Total de slides: {len(prs.slides)}")

    return output_path


if __name__ == "__main__":
    main()
